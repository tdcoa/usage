"test cases for powerpoint templating"
from typing import List
from pathlib import Path
import pytest
import pptx
from tdcsm.pptx import find_placeholders, replace_placeholders


@pytest.fixture(scope="session")
def pptxdir(testdir, tmp_path_factory) -> Path:
	tmpdir = tmp_path_factory.mktemp("tdcsm_pptx")
	for p in (testdir / "test_pptx").iterdir():
		(tmpdir / p.name).write_bytes(p.read_bytes())
	return tmpdir


def phlist(fname: Path) -> List[str]:
	return [str(ph) for ph in find_placeholders(pptx.Presentation(fname))]


def test_ppt_pic(pptxdir: Path) -> None:
	"replace pic placeholder"
	assert phlist(pptxdir / "test_pic.pptx") == [
		'type=pic, data=alldates.png, slide#=1, shape=Rectangle 3'
	]


def test_ppt_col(pptxdir: Path) -> None:
	"replace pic placeholder"
	assert phlist(pptxdir / "test_col.pptx") == [
		'type=col, data=dates.csv, slide#=1, shape=Table 4, loc=0',
		'type=col, data=dates.csv, slide#=1, shape=Table 4, loc=1'
	]


def test_ppt_val(pptxdir: Path) -> None:
	"replace pic placeholder"
	assert phlist(pptxdir / "test_val.pptx") == [
		'type=val, data=birthday.csv, slide#=1, shape=Content Placeholder 1, loc={{val:birthday.csv[1:2]}}'
	]


def test_ppt_all(pptxdir: Path) -> None:
	replace_placeholders(pptxdir / "test_all.pptx", pptxdir)
	ppt = pptx.Presentation(pptxdir / "test_all.pptx")
	shapes = [sh for sn, sl in enumerate(ppt.slides, start=1) for sh in sl.shapes]

	assert shapes[3].image.blob == (pptxdir / "alldates.png").read_bytes()
	assert shapes[4].text == 'Text with one placeholder 1974-12-17'
	assert [[shapes[6].table.cell(r, c).text for c in range(2)] for r in range(3)] == [
		['cal_date', 'item'],
		['2020/01/01', 'Decade'],
		['2000/01/01', 'Millennium']
	]
	assert shapes[7].text == 'TextBox with multiple placehoders: USA received independence on 1776/07/04'
